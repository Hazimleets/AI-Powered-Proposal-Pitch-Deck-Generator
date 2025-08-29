#backend/generator/utils.py

from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor



# -----------------------
# PDF Builder
# -----------------------
def build_pdf(sections, buffer, template="classic"):
    doc = SimpleDocTemplate(buffer)
    styles = getSampleStyleSheet()
    story = []

    if template == "classic":
        # Classic → clean, minimal
        heading_style = styles["Heading2"]
        normal_style = styles["Normal"]

        for section_title, section_content in sections.items():
            story.append(Paragraph(section_title, heading_style))
            story.append(Paragraph(section_content, normal_style))
            story.append(Spacer(1, 12))

    elif template == "modern":
        # Modern → bold headers, blue background
        heading_style = ParagraphStyle(
            'ModernHeading',
            parent=styles['Heading1'],
            textColor=colors.white,
            backColor=colors.HexColor("#1A73E8"),
            fontSize=16,
            spaceAfter=12,
            alignment=1  # centered
        )
        normal_style = ParagraphStyle(
            'ModernNormal',
            parent=styles['Normal'],
            fontName="Helvetica",
            fontSize=12,
            leading=16,
            textColor=colors.HexColor("#333333")
        )

        for section_title, section_content in sections.items():
            story.append(Paragraph(section_title, heading_style))
            story.append(Paragraph(section_content, normal_style))
            story.append(Spacer(1, 20))

    doc.build(story)


# -----------------------
# PPTX Builder
# -----------------------
def build_pptx(sections, file_path, template="classic"):
    prs = Presentation()

    # Define styles
    if template == "classic":
        bg_color = RGBColor(255, 255, 255)   # white
        title_color = RGBColor(0, 51, 102)   # navy blue
        content_color = RGBColor(0, 0, 0)    # black
    else:  # modern
        bg_color = RGBColor(30, 30, 60)      # dark blue
        title_color = RGBColor(255, 255, 255) # white
        content_color = RGBColor(200, 200, 200) # light gray

    # Create slides
    for title, content in sections.items():
        slide_layout = prs.slide_layouts[5]  # Title only
        slide = prs.slides.add_slide(slide_layout)

        # Background color
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = title_color

        # Content box
        left = Pt(50)
        top = Pt(120)
        width = Pt(600)
        height = Pt(350)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        p = tf.add_paragraph()
        p.text = content
        p.font.size = Pt(20)
        p.font.color.rgb = content_color

    prs.save(file_path)

